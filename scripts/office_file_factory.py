"""Erzeugt echte, valide Office-Dateien (xlsx/pptx/docx/pdf) mit Blindtext.

Gedacht fuer Test-Spaces: view-file-Makros und Migrationstools bekommen
Dateien, die echte Previews/Konvertierungen ausloesen. Die Abhaengigkeiten
sind bewusst NICHT Teil der Projekt-Dependencies; Aufruf z. B.:

    uv run --with openpyxl --with python-pptx --with python-docx \
        --with fpdf2 python scripts/office_file_factory.py out/

Alle Generatoren sind deterministisch (random.Random-Instanz uebergeben).
"""

from __future__ import annotations

import io
import random

WORDS = (
    "Antrag Genehmigung Beschaffung Rechnung Freigabe Workflow Personal "
    "Dienstreise Abrechnung Vorlage Formular Prozess Richtlinie Zustaendigkeit "
    "Kostenstelle Budget Vertrag Lieferant Bestellung Inventar Onboarding "
    "Fortbildung Urlaub Vertretung Protokoll Besprechung Projekt Meilenstein"
).split()


def _lorem(rng: random.Random, n: int) -> str:
    return " ".join(rng.choice(WORDS) for _ in range(n))


def make_xlsx(rng: random.Random) -> bytes:
    """Arbeitsmappe mit 2-3 Sheets, Zahlenreihen, Summenformel, Diagramm."""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference

    wb = Workbook()
    for si in range(rng.randint(2, 3)):
        ws = wb.active if si == 0 else wb.create_sheet()
        ws.title = f"{rng.choice(WORDS)}-{si + 1}"
        headers = ["Position", "Menge", "Einzelpreis", "Summe", "Status"]
        ws.append(headers)
        n_rows = rng.randint(30, 200)
        for r in range(2, n_rows + 2):
            menge = rng.randint(1, 50)
            preis = round(rng.uniform(5, 900), 2)
            ws.append(
                [
                    _lorem(rng, 2),
                    menge,
                    preis,
                    f"=B{r}*C{r}",
                    rng.choice(["offen", "genehmigt", "abgelehnt"]),
                ]
            )
        ws.append(["Gesamt", "", "", f"=SUM(D2:D{n_rows + 1})", ""])
        if si == 0:
            chart = BarChart()
            chart.title = "Mengen"
            data = Reference(ws, min_col=2, min_row=1, max_row=min(n_rows, 15))
            chart.add_data(data, titles_from_data=True)
            ws.add_chart(chart, "G2")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx(rng: random.Random) -> bytes:
    """Praesentation mit Titel, Bullet-Folien und einer Tabelle."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = _lorem(rng, 3).title()
    slide.placeholders[1].text = _lorem(rng, 6)

    for _ in range(rng.randint(4, 10)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _lorem(rng, 3).title()
        tf = slide.placeholders[1].text_frame
        tf.text = _lorem(rng, 8)
        for _ in range(rng.randint(2, 5)):
            p = tf.add_paragraph()
            p.text = _lorem(rng, rng.randint(4, 10))
            p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Uebersicht"
    rows, cols = rng.randint(3, 6), 3
    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1.8), Inches(8), Inches(0.8 * rows)
    ).table
    for c, h in enumerate(["Thema", "Verantwortlich", "Termin"]):
        table.cell(0, c).text = h
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = _lorem(rng, 2)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_docx(rng: random.Random) -> bytes:
    """Dokument mit Ueberschriften, Absaetzen und einer Tabelle."""
    from docx import Document

    doc = Document()
    doc.add_heading(_lorem(rng, 3).title(), level=1)
    for _ in range(rng.randint(2, 4)):
        doc.add_heading(_lorem(rng, 2).title(), level=2)
        for _ in range(rng.randint(1, 3)):
            doc.add_paragraph(_lorem(rng, rng.randint(30, 80)) + ".")
    table = doc.add_table(rows=rng.randint(3, 6), cols=3)
    table.style = "Table Grid"
    for row in table.rows:
        for cell in row.cells:
            cell.text = _lorem(rng, 2)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pdf(rng: random.Random) -> bytes:
    """PDF mit 1-3 Seiten Blindtext."""
    from fpdf import FPDF

    pdf = FPDF()
    for _ in range(rng.randint(1, 3)):
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, _lorem(rng, 3).title(), new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", size=11)
        for _ in range(rng.randint(3, 6)):
            pdf.multi_cell(0, 6, _lorem(rng, rng.randint(40, 90)) + ".")
            pdf.ln(3)
    return bytes(pdf.output())


MAKERS = {
    "xlsx": make_xlsx,
    "pptx": make_pptx,
    "docx": make_docx,
    "pdf": make_pdf,
}


def make_file(rng: random.Random, ext: str) -> bytes:
    """Erzeugt eine Datei passend zur Endung (KeyError bei unbekannter)."""
    return MAKERS[ext](rng)


if __name__ == "__main__":
    import pathlib
    import sys

    out = pathlib.Path(sys.argv[1] if len(sys.argv) > 1 else ".")
    out.mkdir(parents=True, exist_ok=True)
    rng = random.Random(42)
    for ext in MAKERS:
        p = out / f"beispiel.{ext}"
        p.write_bytes(make_file(rng, ext))
        print(p, p.stat().st_size, "Bytes")
